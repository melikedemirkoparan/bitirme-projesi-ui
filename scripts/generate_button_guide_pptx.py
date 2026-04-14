"""
Generate a Turkish PowerPoint (.pptx) presentation that walks through the
Patent Drafting Tool UI, button by button, and explains the technical
architecture for the new Inventor_QA feature.

Run from the repo root:
    python scripts/generate_button_guide_pptx.py

Output: docs/Patent_Drafting_Tool_Sunum.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "Patent_Drafting_Tool_Sunum.pptx"

SCREENSHOTS_DIR = Path(
    r"C:\Users\melike\OneDrive - CC Bilgi Teknolojileri\Resimler\Ekran görüntüleri"
)
SCREENSHOTS = {
    "dashboard": SCREENSHOTS_DIR / "Ekran görüntüsü 2026-04-14 011914.png",
    "workspace_full": SCREENSHOTS_DIR / "Ekran görüntüsü 2026-04-14 011929.png",
    "workspace_alt": SCREENSHOTS_DIR / "Ekran görüntüsü 2026-04-14 012018.png",
    "draft_editor": SCREENSHOTS_DIR / "Ekran görüntüsü 2026-04-14 012112.png",
    "upload_excel": SCREENSHOTS_DIR / "Ekran görüntüsü 2026-04-14 012052.png",
    "extract_modal": SCREENSHOTS_DIR / "Ekran görüntüsü 2026-04-14 012036.png",
    "project_settings": SCREENSHOTS_DIR / "Ekran görüntüsü 2026-04-14 012311.png",
}

# Colors (dark theme that matches the app)
BG_DARK = RGBColor(0x0A, 0x0F, 0x1A)
BG_PANEL = RGBColor(0x1A, 0x22, 0x36)
BG_HEADER = RGBColor(0x1E, 0x3C, 0x82)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
TEXT_PRIMARY = RGBColor(0xE8, 0xEA, 0xF0)
TEXT_SECONDARY = RGBColor(0x88, 0x99, 0xB4)
TEXT_MUTED = RGBColor(0x5A, 0x6A, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARNING = RGBColor(0xEA, 0xB3, 0x08)
SUCCESS = RGBColor(0x22, 0xC5, 0x5E)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK_LAYOUT = prs.slide_layouts[6]


def add_blank_slide():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, italic=False, color=TEXT_PRIMARY,
                align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=14, color=TEXT_PRIMARY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        if isinstance(item, tuple):
            label, body = item
            r1 = p.add_run()
            r1.text = "• " + label + " "
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = ACCENT
            r2 = p.add_run()
            r2.text = body
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = "• " + item
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def slide_header(slide, title_text, subtitle_text=None):
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG_HEADER
    bar.line.fill.background()
    add_textbox(
        slide, Inches(0.5), Inches(0.18), SLIDE_W - Inches(1.0), Inches(0.55),
        title_text, size=26, bold=True, color=WHITE,
    )
    if subtitle_text:
        add_textbox(
            slide, Inches(0.5), Inches(0.95), SLIDE_W - Inches(1.0), Inches(0.4),
            subtitle_text, size=14, italic=True, color=TEXT_SECONDARY,
        )


def add_image(slide, key, left, top, max_width, max_height):
    path = SCREENSHOTS.get(key)
    if not path or not path.exists():
        return
    pic = slide.shapes.add_picture(str(path), left, top, width=max_width)
    if pic.height > max_height:
        ratio = max_height / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = max_height
        pic.left = int(left + (max_width - pic.width) / 2)
    pic.line.color.rgb = ACCENT
    pic.line.width = Pt(1)


# ─────────────────────────────────────────────────────────────────
# 1) COVER
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()

title_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SLIDE_W, Inches(2.7)
)
title_box.fill.solid()
title_box.fill.fore_color.rgb = BG_PANEL
title_box.line.fill.background()

add_textbox(
    slide, Inches(0.5), Inches(2.7), SLIDE_W - Inches(1.0), Inches(0.9),
    "Patent Drafting Tool", size=44, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER,
)
add_textbox(
    slide, Inches(0.5), Inches(3.6), SLIDE_W - Inches(1.0), Inches(0.6),
    "Buton Rehberi & Teknik Sunum",
    size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
)
add_textbox(
    slide, Inches(0.5), Inches(4.3), SLIDE_W - Inches(1.0), Inches(0.5),
    "Arayüzdeki tüm butonların ne işe yaradığı + Inventor_QA modülünün teknik mimarisi",
    size=14, italic=True, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER,
)
add_textbox(
    slide, Inches(0.5), Inches(6.6), SLIDE_W - Inches(1.0), Inches(0.4),
    "Sürüm: feature/inventor-qa-section",
    size=11, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
)


# ─────────────────────────────────────────────────────────────────
# 2) AGENDA
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "İçindekiler", "Sunumun akışı")
add_bullets(
    slide, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
    [
        "1. Ana Sayfa (Dashboard) — proje listesi ve butonları",
        "2. Workspace — Üst Navbar",
        "3. Workspace — Sol Panel: Claims & Elements",
        "4. Workspace — Orta Panel: LLM Extraction Queue",
        "5. Workspace — Sağ Panel: Claim Draft Editor",
        "6. Modal: Upload Excel Data",
        "7. Modal: Extract Elements from BBF + Report (+ Inventor_QA)",
        "8. Modal: Project Settings",
        "9. Hızlı İş Akışı (sıfırdan draft üretme)",
        "10. Teknik Bilgi — Kod ve Frontend Mimarisi",
        "11. Yeni Inventor_QA Modülü — Backend & Frontend Detayı",
    ],
    size=16,
)


# ─────────────────────────────────────────────────────────────────
# Helper to build a "screen + buttons" 2-column slide
# ─────────────────────────────────────────────────────────────────
def screen_slide(title, subtitle, screenshot_key, buttons, image_left=Inches(0.5),
                 image_max_w=Inches(6.4), image_max_h=Inches(5.6),
                 buttons_left=Inches(7.2), buttons_size=12):
    slide = add_blank_slide()
    slide_header(slide, title, subtitle)
    add_image(slide, screenshot_key, image_left, Inches(1.5),
              image_max_w, image_max_h)
    add_bullets(
        slide, buttons_left, Inches(1.5), Inches(5.6), Inches(5.6),
        buttons, size=buttons_size,
    )


# ─────────────────────────────────────────────────────────────────
# 3) ANA SAYFA (Dashboard)
# ─────────────────────────────────────────────────────────────────
screen_slide(
    "1. Ana Sayfa (Dashboard)",
    "Proje listesi + tanıtım panosu + üst navbar",
    "dashboard",
    [
        ("Data Loaded / No Data:", "Excel verisinin yüklü olup olmadığını gösterir."),
        ("Back to Home:", "Workspace'ten ana sayfaya döner."),
        ("Upload Data:", "Excel yükleme modalını açar."),
        ("Upload Excel Data (sol):", "Aynı yükleme modalı, kolaylık için sol panelde de var."),
        ("+ New Project:", "Yeni proje oluşturma modalı (ad + owner)."),
        ("Inputs (proje kartında):", "Patent Inputs modalı — Inventor_QA metni ve dökümanları."),
        ("Edit (proje kartında):", "Proje adı ve owner düzenleme."),
        ("Open / Karta tıklama:", "Projeyi açıp workspace ekranına geçer."),
    ],
)


# ─────────────────────────────────────────────────────────────────
# 4) WORKSPACE — ÜST NAVBAR
# ─────────────────────────────────────────────────────────────────
screen_slide(
    "2. Workspace — Üst Navbar",
    "Açık projenin üst çubuğu",
    "workspace_full",
    [
        ("Proje Adı:", "Açık olan projenin adı; yanında owner görünür."),
        ("✎ Rename:", "Projenin adını/sahibini değiştirir (Edit Project modalı)."),
        ("Back to Home:", "Ana sayfaya geri döner."),
        ("Upload Data:", "Workspace içinden de Excel yüklenebilir."),
    ],
)


# ─────────────────────────────────────────────────────────────────
# 5) SOL PANEL — CLAIMS & ELEMENTS
# ─────────────────────────────────────────────────────────────────
screen_slide(
    "3. Sol Panel — Claims & Elements",
    "Patent istemleri ve element ağacı",
    "workspace_alt",
    [
        ("+ Add Claim:", "Yeni istem oluşturur (independent/dependent + apparatus/method)."),
        ("Claim Numarası:", "Tıklayınca o istemi seçer; Ready/Incomplete rozeti var."),
        ("Genişlet/Daralt oku:", "Element ağacını açıp kapatır."),
        ("+ Add Elements:", "Patent havuzundan element bağlar."),
        ("Yukarı/Aşağı oku:", "Element sırasını değiştirir."),
        ("Yenile:", "Element ağacını backend'den tazeler."),
        ("Search elements:", "İsme göre filtreler."),
        ("Element satırı:", "Tıklayınca tanım modalı açılır."),
        ("Çöp ikonu (claim/element):", "Claim'i siler veya element bağını çözer."),
    ],
    buttons_size=11,
)


# ─────────────────────────────────────────────────────────────────
# 6) ORTA PANEL — LLM EXTRACTION QUEUE
# ─────────────────────────────────────────────────────────────────
screen_slide(
    "4. Orta Panel — LLM Extraction Queue",
    "Patent element havuzu ve otomatik çıkarım",
    "workspace_full",
    [
        ("⚡ Extract:", "BBF + Report (+ opsiyonel Inventor_QA) yükleyip otomatik element çıkarır."),
        ("+ Add Element:", "Manuel olarak yeni element oluşturur."),
        ("DRAG (sürükleme):", "Element'i bir claim kartına sürükleyerek bağlar."),
        ("Edit:", "Element Definition modalını açar (ad, ref, definition, AI suggest)."),
        ("Çöp ikonu:", "Element'i havuzdan tamamen siler."),
    ],
)


# ─────────────────────────────────────────────────────────────────
# 7) SAĞ PANEL — CLAIM DRAFT EDITOR
# ─────────────────────────────────────────────────────────────────
screen_slide(
    "5. Sağ Panel — Claim Draft Editor",
    "Seçili istem için draft yazma alanı",
    "draft_editor",
    [
        ("⚙ Settings:", "Project Settings modalı (Invention Context, BBF Text, LLM URL)."),
        ("Assemble with Report:", "Tüm tanımlardan bir rapor üretip draft alanına yazar."),
        ("Save:", "Draft metnini PostgreSQL'e kaydeder; '✓ Saved' geri bildirimi gösterir."),
        ("Go to AI Draft:", "Composer ekranına geçip AI ile tam draft üretir."),
        ("Draft Composer:", "Direkt Draft Composer sayfasına geçiş yapar."),
    ],
)


# ─────────────────────────────────────────────────────────────────
# 8) MODAL — UPLOAD EXCEL DATA
# ─────────────────────────────────────────────────────────────────
screen_slide(
    "6. Modal — Upload Excel Data",
    "RAG vektör veritabanını besleyen Excel yükleme",
    "upload_excel",
    [
        ("Click to select Excel file:", "Bilgisayardan .xlsx dosyası seçer."),
        ("Embedding süreci:", "Arka planda çalışır; durum yeşil rozetle gösterilir (örn. 'Loaded 9133 docs. FAISS index ready.')."),
        ("Close:", "Modalı kapatır; arka plan işlemi etkilenmez."),
    ],
    image_max_w=Inches(6.4),
    image_max_h=Inches(5.5),
)


# ─────────────────────────────────────────────────────────────────
# 9) MODAL — EXTRACT ELEMENTS (with Inventor_QA)
# ─────────────────────────────────────────────────────────────────
screen_slide(
    "7. Modal — Extract Elements from BBF + Report",
    "BBF + Report (+ Inventor_QA) ile otomatik element çıkarımı",
    "extract_modal",
    [
        ("BBF Document:", "Buluş Bildirim Formu dökümanı (zorunlu)."),
        ("Report Document:", "Araştırma raporu dökümanı (zorunlu)."),
        ("Buluşçu ile Yazışmalar (Inventor_QA):", "Opsiyonel; dosya seçer seçmez otomatik yüklenir ve projede saklanır."),
        ("Cancel:", "Modalı kapatır."),
        ("⚡ Extract Elements:", "Pipeline'ı çalıştırır; çıkarılan element'ler havuza eklenir."),
    ],
    image_max_w=Inches(5.5),
    image_max_h=Inches(5.6),
    buttons_size=12,
)


# ─────────────────────────────────────────────────────────────────
# 10) MODAL — PROJECT SETTINGS
# ─────────────────────────────────────────────────────────────────
screen_slide(
    "8. Modal — Project Settings",
    "AI üretimi için bağlam ve LLM ayarları",
    "project_settings",
    [
        ("Invention Context:", "Buluşu özetleyen tek satırlık context cümlesi."),
        ("BBF Text:", "BBF tam metni; AI üretim modüllerine girdi olur."),
        ("LLM API URL (Colab ngrok):", "Colab'da çalışan Mistral-7B'nin uzaktan erişim adresi."),
        ("Cancel:", "Kaydetmeden kapatır."),
        ("Save:", "Ayarları uygulamanın cache'ine yazar."),
    ],
    image_max_w=Inches(5.0),
    image_max_h=Inches(5.6),
)


# ─────────────────────────────────────────────────────────────────
# 11) HIZLI İŞ AKIŞI
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "9. Hızlı İş Akışı", "Sıfırdan tam patent draft'ı için 10 adım")
add_bullets(
    slide, Inches(0.7), Inches(1.4), Inches(12), Inches(5.7),
    [
        "1. Excel verisini yükleyin (Upload Data) → 'Data Loaded' rozeti yeşile dönmeli.",
        "2. + New Project ile proje oluşturun (ad + owner).",
        "3. (Opsiyonel) Inputs butonu ile Inventor_QA metni / dökümanı ekleyin.",
        "4. Open ile projeyi açın → workspace ekranına geçin.",
        "5. ⚡ Extract ile BBF + Report (+ Inventor_QA) yükleyip element'leri çıkarın.",
        "6. + Add Claim ile bir veya birden fazla istem oluşturun.",
        "7. Element'leri claim'lere bağlayın (+ Add Elements veya DRAG).",
        "8. Element Definition modallarını açıp tanımları doldurun (gerekirse AI Suggest).",
        "9. Bir claim seçin → draft metnini yazıp Save'e basın.",
        "10. Go to AI Draft → Generate Draft ile AI çıktısı üretin.",
    ],
    size=16,
)


# ─────────────────────────────────────────────────────────────────
# 12) TEKNOLOJİ STACK
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "10. Kullanılan Teknolojiler", "Stack — neyin nerede çalıştığı")
add_bullets(
    slide, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
    [
        ("Backend:", "FastAPI (Python 3.12) — REST API ve statik dosya sunumu."),
        ("Veritabanı:", "PostgreSQL 16 — projeler, claim'ler, element'ler, inventor_qa, dökümanlar."),
        ("ORM:", "SQLAlchemy 2.x (Mapped tipler) + Alembic (şema migration'ları)."),
        ("Vektör veritabanı:", "ChromaDB / FAISS + sentence-transformers (multilingual-e5-base)."),
        ("AI / LLM:", "RAG akışı + Colab üzerinde Mistral-7B (ngrok ile uzaktan)."),
        ("Frontend:", "Vanilla HTML + CSS + JavaScript — build adımı yok, SPA gibi davranır."),
        ("Container:", "Docker Compose ile iki servis: db (postgres) ve app (FastAPI)."),
        ("Versiyon kontrol:", "Git + GitHub (feature branch akışı)."),
    ],
    size=15,
)


# ─────────────────────────────────────────────────────────────────
# 13) KLASÖR YAPISI
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "11. Klasör Yapısı", "Önemli dizinler ve sorumlulukları")
tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.7))
tf = tb.text_frame
tf.word_wrap = True
code = (
    "paten_draft_backend/\n"
    "├── app/\n"
    "│   ├── main.py                FastAPI + router bağlama\n"
    "│   ├── database.py            SQLAlchemy session\n"
    "│   ├── config.py              Ortam ayarları\n"
    "│   ├── models/                ORM (Patent, Claim, Element, InventorQA…)\n"
    "│   ├── schemas/               Pydantic istek/cevap şemaları\n"
    "│   ├── services/              İş mantığı (CRUD, dosya işlemleri)\n"
    "│   ├── routes/                FastAPI endpoint'leri\n"
    "│   ├── ingestion/             Excel → ChromaDB embedding üretimi\n"
    "│   └── rag_engine.py          RAG sorgu motoru\n"
    "├── alembic/versions/          DB migration dosyaları\n"
    "├── static/\n"
    "│   ├── home.html              Tek sayfa SPA şablonu\n"
    "│   ├── jss/app.js             Tüm frontend mantığı\n"
    "│   └── cs/style.css           Koyu tema\n"
    "├── docker-compose.yml         db + app servisleri\n"
    "└── Dockerfile"
)
p = tf.paragraphs[0]
r = p.add_run()
r.text = code
r.font.name = "Consolas"
r.font.size = Pt(13)
r.font.color.rgb = TEXT_PRIMARY


# ─────────────────────────────────────────────────────────────────
# 14) FRONTEND MİMARİSİ
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "12. Frontend Mimarisi", "static/home.html + jss/app.js + cs/style.css")
add_bullets(
    slide, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
    [
        "Tek 'home.html' içinde birden fazla 'page' tanımlı: page-dashboard, page-editor, page-composer.",
        "Sayfa geçişleri showPage(id) ile yapılır; sadece bir sayfa 'active' class'ına sahip olur.",
        "Tüm modal'lar aynı dosyada tanımlı; 'modal-overlay' class'ına 'active' eklenince açılır.",
        "State global JS değişkenleriyle yönetilir: _projects, _claims, _elements, _patentId vb.",
        "API çağrıları için 'api(path, opts)' helper'ı — otomatik /api prefix + JSON parse.",
        "Cache busting: app.js?v=5 ve style.css?v=5 ile tarayıcı eski sürümü değil yenisini indirir.",
    ],
    size=15,
)


# ─────────────────────────────────────────────────────────────────
# 15) BACKEND MİMARİSİ
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "13. Backend Mimarisi", "4 katman: model → schema → service → route")
add_bullets(
    slide, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5),
    [
        ("Model (app/models/):", "SQLAlchemy ORM — DB tablolarının Python karşılığı."),
        ("Schema (app/schemas/):", "Pydantic — istek body ve cevap formatı, doğrulama."),
        ("Service (app/services/):", "İş mantığı — service'ler DB ile konuşur, route'lar service'i çağırır."),
        ("Route (app/routes/):", "FastAPI APIRouter — HTTP endpoint tanımları."),
        ("Bağlantı:", "main.py routerları include_router ile uygulamaya ekler."),
        ("DB session:", "Depends(get_db) ile her isteğe enjekte edilir; istek sonunda kapanır."),
        ("Uzun işlemler:", "Excel embedding ve BBF pipeline thread'lerde arka planda çalışır; status endpoint'i ile poll edilir."),
    ],
    size=14,
)


# ─────────────────────────────────────────────────────────────────
# 16) DB TABLOLARI
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "14. Veritabanı Şeması", "PostgreSQL ana tabloları")
add_bullets(
    slide, Inches(0.7), Inches(1.4), Inches(12), Inches(5.8),
    [
        ("patent:", "patent_id, patent_name, patent_owner, patent_draft, timestamps"),
        ("claim:", "claim_id, patent_id, claim_number, dependency_type, category, parent_claim_id, claim_text"),
        ("element:", "element_id, patent_id, element_name, reference_number, definition_text"),
        ("claim_element:", "claim_id + element_id (many-to-many) + order_index"),
        ("invention_disclosure:", "patent_id (1:1), prior_art_and_problems, closest_prior_patents, novel_features"),
        ("research_report:", "patent_id (1:1), executive_summary, search_strategy, classification_and_keywords, element_patent_analysis"),
        ("inventor_qa:", "patent_id (1:1), questions_and_answers (text)"),
        ("inventor_qa_document  [YENİ]:", "qna_id (FK), original_filename, stored_filename, mime_type, size_bytes, created_at"),
    ],
    size=13,
)


# ─────────────────────────────────────────────────────────────────
# 17) YENİ INVENTOR_QA — BACKEND
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(
    slide,
    "15. Yeni Inventor_QA — Backend",
    "Buluşçu ile Yazışmalar modülünün backend tarafı",
)
add_bullets(
    slide, Inches(0.7), Inches(1.4), Inches(12), Inches(5.8),
    [
        ("Model:", "app/models/inventor_qa.py içine InventorQADocument eklendi (cascade silme)."),
        ("Schema:", "app/schemas/inventor_qa.py — InventorQARead, InventorQADocumentRead, InventorQAUpdate."),
        ("Service:", "app/services/inventor_qa_service.py — get/upsert text + add/get/list/delete document."),
        ("Route:", "app/routes/inventor_qa.py — /api/patents/{patent_id}/inventor-qa altında 5 endpoint."),
        ("Migration:", "alembic/versions/7c0e9b41ad22_add_inventor_qa_document.py."),
        ("Dosya saklama:", "settings.uploads_path/inventor_qa/{patent_id}/ klasörüne yazılır; metadata DB'de tutulur."),
        ("App'e bağlama:", "main.py içinde app.include_router(inventor_qa_router) eklendi."),
    ],
    size=14,
)


# ─────────────────────────────────────────────────────────────────
# 18) YENİ INVENTOR_QA — API ENDPOINT'LERİ
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "16. Yeni API Endpoint'leri", "RESTful arayüz")
tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True
code = (
    "GET    /api/patents/{patent_id}/inventor-qa\n"
    "       → Q&A metnini ve yüklenmiş doküman listesini döndürür\n\n"
    "PUT    /api/patents/{patent_id}/inventor-qa\n"
    "       → Q&A metnini günceller (body: { questions_and_answers })\n\n"
    "POST   /api/patents/{patent_id}/inventor-qa/documents\n"
    "       → multipart/form-data ile doküman yükler\n\n"
    "GET    /api/patents/{patent_id}/inventor-qa/documents/{document_id}\n"
    "       → Yüklenmiş dokümanı indirir (FileResponse)\n\n"
    "DELETE /api/patents/{patent_id}/inventor-qa/documents/{document_id}\n"
    "       → Dokümanı diskten ve DB'den siler"
)
p = tf.paragraphs[0]
r = p.add_run()
r.text = code
r.font.name = "Consolas"
r.font.size = Pt(13)
r.font.color.rgb = TEXT_PRIMARY


# ─────────────────────────────────────────────────────────────────
# 19) YENİ INVENTOR_QA — FRONTEND
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "17. Yeni Inventor_QA — Frontend", "static/home.html + static/jss/app.js")
add_bullets(
    slide, Inches(0.7), Inches(1.4), Inches(12), Inches(5.8),
    [
        ("Patent Inputs Modalı:", "home.html'e eklendi — Inventor_QA kartı (textarea + dosya yükleme + döküman listesi)."),
        ("Proje kartı butonları:", "Her proje kartına 📋 Inputs ve ✎ Edit butonları eklendi (sıralı, dar sidebar'a sığar)."),
        ("Workspace Rename:", "Üst navbar'a ✎ Rename butonu eklendi; aynı edit modalını açar."),
        ("Extract modalı:", "BBF + Report'a ek 3. yükleme alanı (Buluşçu ile Yazışmalar) eklendi; dosya seçer seçmez otomatik upload yapar."),
        ("JS fonksiyonları:", "openPatentInputs, loadInventorQa, uploadInventorQaDocument, deleteInventorQaDocument, saveInventorQaText, onInventorQaExtractFileChange."),
        ("CSS:", "patent-input-card, patent-input-doc-row, project-item--stacked stilleri."),
        ("Cache busting:", "Tüm bu değişiklikler app.js?v=5 / style.css?v=5 ile garanti edilir."),
    ],
    size=13,
)


# ─────────────────────────────────────────────────────────────────
# 20) İSTEK-CEVAP AKIŞI (örnek)
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "18. İstek-Cevap Akışı", "Örnek: Q&A metnini kaydetme")
add_bullets(
    slide, Inches(0.7), Inches(1.4), Inches(12), Inches(5.8),
    [
        "1. Kullanıcı 📋 Inputs butonuna tıklar → frontend openPatentInputs() çağrılır.",
        "2. Modal açılır → loadInventorQa() ile GET /api/patents/{id}/inventor-qa çağrılır.",
        "3. Backend route → service → SQLAlchemy SELECT → mevcut kayıt döner.",
        "4. Kullanıcı textarea'ya yazıp 'Save Q&A Text' butonuna basar → saveInventorQaText().",
        "5. Frontend PUT /api/patents/{id}/inventor-qa endpoint'ine JSON gönderir.",
        "6. FastAPI route → service → SQLAlchemy UPDATE → güncellenmiş kayıt döner.",
        "7. Frontend cevabı alır, doküman listesini yeniler ve 'Saved' uyarısı gösterir.",
    ],
    size=15,
)


# ─────────────────────────────────────────────────────────────────
# 21) ÇALIŞTIRMA
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()
slide_header(slide, "19. Çalıştırma & Test", "Docker Compose komutları")
tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True
code = (
    "# Tek komutla başlat\n"
    "docker compose up -d --build\n\n"
    "# Migration container açılışında otomatik çalışır (start.sh)\n"
    "# alembic upgrade head\n\n"
    "# Loglar\n"
    "docker compose logs -f app\n\n"
    "# Sağlık kontrolü\n"
    "curl http://localhost:8000/api/health\n\n"
    "# Tarayıcı\n"
    "http://localhost:8000/\n\n"
    "# Branch ve push\n"
    "git checkout -b feature/inventor-qa-section\n"
    "git push melike feature/inventor-qa-section\n"
    "git push origin feature/inventor-qa-section"
)
p = tf.paragraphs[0]
r = p.add_run()
r.text = code
r.font.name = "Consolas"
r.font.size = Pt(13)
r.font.color.rgb = TEXT_PRIMARY


# ─────────────────────────────────────────────────────────────────
# 22) KAPANIŞ
# ─────────────────────────────────────────────────────────────────
slide = add_blank_slide()

box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, Inches(2.7), SLIDE_W, Inches(2.2)
)
box.fill.solid()
box.fill.fore_color.rgb = BG_PANEL
box.line.fill.background()

add_textbox(
    slide, Inches(0.5), Inches(2.95), SLIDE_W - Inches(1.0), Inches(0.9),
    "Teşekkürler!", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
)
add_textbox(
    slide, Inches(0.5), Inches(3.95), SLIDE_W - Inches(1.0), Inches(0.6),
    "Sorularınızı bekliyorum.",
    size=18, italic=True, color=ACCENT, align=PP_ALIGN.CENTER,
)
add_textbox(
    slide, Inches(0.5), Inches(5.4), SLIDE_W - Inches(1.0), Inches(0.5),
    "GitHub: melikedemirkoparan/bitirme-projesi-ui (branch: feature/inventor-qa-section)",
    size=12, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
)


OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT))
print(f"Wrote {OUTPUT}")
